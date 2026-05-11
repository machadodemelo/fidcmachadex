from __future__ import annotations

from datetime import date
from io import BytesIO
import json
from pathlib import Path
import tempfile
from types import SimpleNamespace
import unittest
import zipfile

from openpyxl import load_workbook
import pandas as pd

from services.mercado_livre_dashboard import (
    build_consolidated_monthly_base,
    build_consolidated_snapshot_excel_bytes,
    build_excel_export_bytes,
    build_fund_monthly_base,
    build_mercado_livre_outputs,
    build_wide_table,
    load_outputs_from_cache,
    MercadoLivreOutputs,
    order_period_columns_desc,
    portfolio_identity_key,
    save_outputs_to_cache,
)
from services.mercado_livre_ppt_export import build_pptx_export_bytes
from services.portfolio_store import PortfolioFund, PortfolioRecord
from services.mercado_livre_visuals import npl_coverage_chart, pl_subordination_chart
from tabs.tab_mercado_livre import (
    _MERCADO_LIVRE_UI_CSS,
    _build_mercado_livre_guide_markdown,
    _dense_wide_value,
    _display_window_bounds,
    _display_window_months,
    _build_credit_monitor_for_display,
    _period_with_yoy_lookback,
    _tag_outputs_requested_period,
    _filter_outputs_by_competencia_months,
    _render_wide_table_html,
    _resolve_existing_portfolio_for_save,
)
from services.ime_period import build_custom_period


class MercadoLivreDashboardTests(unittest.TestCase):
    def test_build_fund_monthly_base_calculates_accumulated_npl_and_ex360(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={
                1: 100.0,
                2: 50.0,
                3: 30.0,
                4: 20.0,
                5: 10.0,
                6: 5.0,
                7: 15.0,
                8: 7.0,
                9: 3.0,
                10: 0.0,
            },
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        row = monthly.iloc[0]

        self.assertAlmostEqual(25.0, row["subordinacao_total_pct"])
        self.assertAlmostEqual(60.0, row["npl_over90"])
        self.assertAlmostEqual(10.0, row["npl_over360"])
        self.assertAlmostEqual(50.0, row["npl_over90_ex360"])
        self.assertAlmostEqual(990.0, row["carteira_ex360"])
        self.assertAlmostEqual(50.0 / 990.0 * 100.0, row["npl_over90_ex360_pct"], places=6)
        self.assertTrue(bool(row["pdd_ex360_calculavel"]))
        self.assertAlmostEqual(10.0, row["baixa_over360_carteira"])
        self.assertAlmostEqual(10.0, row["baixa_over360_pdd"])
        self.assertAlmostEqual(0.0, row["baixa_over360_pl"])
        self.assertAlmostEqual(90.0, row["pdd_ex360"])
        self.assertAlmostEqual(900.0, row["carteira_liquida_ex360"])
        self.assertNotIn("PDD ex-360 não calculável", row["warnings"])

    def test_ex360_writeoff_reduces_pl_when_pdd_is_insufficient(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=5.0,
            buckets={8: 10.0},
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        row = monthly.iloc[0]

        self.assertAlmostEqual(10.0, row["baixa_over360_carteira"])
        self.assertAlmostEqual(5.0, row["baixa_over360_pdd"])
        self.assertAlmostEqual(5.0, row["baixa_over360_pl"])
        self.assertAlmostEqual(0.0, row["pdd_ex360"])
        self.assertAlmostEqual(95.0, row["pl_total_ex360"])
        self.assertAlmostEqual(20.0, row["pl_subordinada_mezz_ex360"])
        self.assertIn("PDD menor que Over 360", row["warnings"])

    def test_pl_chart_uses_ex360_subordination_series(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=5.0,
            buckets={8: 10.0},
        )
        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)

        chart_payload = json.dumps(pl_subordination_chart(monthly).to_dict(), ensure_ascii=False)

        self.assertIn("Subordinada + Mez ex-360", chart_payload)
        self.assertIn("% Subordinação Total ex-360", chart_payload)
        self.assertIn("21,1%", chart_payload)
        for legacy_color in ("#1f77b4", "#8ecae6", "#d1495b", "#1b998b"):
            self.assertNotIn(legacy_color, chart_payload)

    def test_npl_chart_uses_meli_palette_and_last_point_labels(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 8: 10.0},
        )
        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)

        chart_payload = json.dumps(npl_coverage_chart(monthly).to_dict(), ensure_ascii=False)

        self.assertIn("#000000", chart_payload)
        self.assertIn("#E47811", chart_payload)
        self.assertIn("Séries", chart_payload)
        self.assertIn("NPL Over 90d ex-360 / Carteira", chart_payload)
        self.assertIn("PDD Ex / NPL Over 90d ex-360", chart_payload)
        self.assertIn("4,0%", chart_payload)
        self.assertIn("225,0%", chart_payload)
        for legacy_color in ("#1f77b4", "#8ecae6", "#d1495b", "#1b998b"):
            self.assertNotIn(legacy_color, chart_payload)

    def test_official_pl_overrides_class_sum_and_keeps_reconciliation(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        official_pl = pd.DataFrame(
            [{"competencia": "01/2026", "pl_total_oficial": 110.0, "pl_total_oficial_source_status": "reported_value"}]
        )

        monthly = build_fund_monthly_base(
            cnpj="11111111000111",
            fund_name="FIDC A",
            dashboard=dashboard,
            official_pl_history_df=official_pl,
        )
        row = monthly.iloc[0]

        self.assertAlmostEqual(110.0, row["pl_total"])
        self.assertAlmostEqual(100.0, row["pl_total_classes"])
        self.assertAlmostEqual(10.0, row["pl_reconciliacao_delta"])
        self.assertEqual("PATRLIQ/VL_PATRIM_LIQ", row["pl_total_usado_fonte"])
        self.assertIn("PL oficial diverge", row["warnings"])

    def test_consolidated_base_sums_absolute_values_and_recalculates_ratios(self) -> None:
        dashboard_a = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard_b = _dashboard(
            fund_name="FIDC B",
            cnpj="22222222000122",
            pl_total=200.0,
            pl_senior=160.0,
            pl_mezz=0.0,
            pl_sub=40.0,
            carteira=2_000.0,
            pdd=200.0,
            buckets={4: 80.0, 7: 20.0, 8: 0.0},
        )
        monthly_a = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard_a)
        monthly_b = build_fund_monthly_base(cnpj="22222222000122", fund_name="FIDC B", dashboard=dashboard_b)

        consolidated = build_consolidated_monthly_base(
            portfolio_name="Carteira",
            fund_monthly_frames={"11111111000111": monthly_a, "22222222000122": monthly_b},
        )
        row = consolidated.iloc[0]

        self.assertAlmostEqual(300.0, row["pl_total"])
        self.assertAlmostEqual(65.0, row["pl_subordinada_mezz"])
        self.assertAlmostEqual(65.0 / 300.0 * 100.0, row["subordinacao_total_pct"], places=6)
        self.assertAlmostEqual(160.0, row["npl_over90"])
        self.assertAlmostEqual(300.0 / 160.0 * 100.0, row["pdd_npl_over90_pct"], places=6)

    def test_roll_rate_uses_current_bucket_over_previous_month_base(self) -> None:
        competencias = ["01/2026", "02/2026"]
        dashboard = SimpleNamespace(
            competencias=competencias,
            fund_info={"nome_fundo": "FIDC A", "cnpj_fundo": "11111111000111", "nome_classe": "Classe única"},
            subordination_history_df=pd.DataFrame(
                [
                    {
                        "competencia": competencia,
                        "competencia_dt": pd.Timestamp(f"2026-{month:02d}-01"),
                        "pl_total": 100.0,
                        "pl_senior": 75.0,
                        "pl_mezzanino": 15.0,
                        "pl_subordinada_strict": 10.0,
                        "pl_subordinada": 25.0,
                    }
                    for month, competencia in enumerate(competencias, start=1)
                ]
            ),
            default_history_df=pd.DataFrame(
                [
                    {
                        "competencia": competencia,
                        "competencia_dt": pd.Timestamp(f"2026-{month:02d}-01"),
                        "direitos_creditorios": 1_000.0,
                        "direitos_creditorios_fonte": "teste",
                        "provisao_total": 100.0,
                    }
                    for month, competencia in enumerate(competencias, start=1)
                ]
            ),
            dc_canonical_history_df=pd.DataFrame(),
            default_buckets_history_df=pd.DataFrame(
                [
                    {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "ordem": 1, "faixa": "Até 30 dias", "valor": 100.0},
                    {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "ordem": 2, "faixa": "31 a 60 dias", "valor": 20.0},
                    {"competencia": "02/2026", "competencia_dt": pd.Timestamp("2026-02-01"), "ordem": 1, "faixa": "Até 30 dias", "valor": 80.0},
                    {"competencia": "02/2026", "competencia_dt": pd.Timestamp("2026-02-01"), "ordem": 2, "faixa": "31 a 60 dias", "valor": 49.0},
                ]
            ),
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        jan = monthly.loc[monthly["competencia"] == "01/2026"].iloc[0]
        feb = monthly.loc[monthly["competencia"] == "02/2026"].iloc[0]

        self.assertAlmostEqual(980.0, jan["carteira_em_dia_mais_ate30"])
        self.assertTrue(pd.isna(jan["roll_rate_base_t_minus_1"]))
        self.assertAlmostEqual(980.0, feb["roll_rate_base_t_minus_1"])
        self.assertAlmostEqual(5.0, feb["roll_rate_31_60_pct"])

        wide = build_wide_table(monthly, scope_name="FIDC A")
        base_row = wide.loc[wide["Métrica"] == "Carteira em dia + atrasada até 30d (t-1)"].iloc[0]
        roll_row = wide.loc[wide["Métrica"] == "Roll Rate"].iloc[0]
        self.assertEqual("R$ 980,00", base_row["fev/26"])
        self.assertEqual("N/D", base_row["jan/26"])
        self.assertEqual("5,00%", roll_row["fev/26"])

    def test_wide_table_and_excel_export_include_required_blocks(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        wide = build_wide_table(outputs.fund_monthly["11111111000111"], scope_name="FIDC A")

        self.assertIn("jan/26", wide.columns)
        self.assertIn("7. Visão Ex-Vencidos > 360d", set(wide["Bloco"]))
        self.assertIn("NPL Over 90d", set(wide["Métrica"]))

        excel_bytes = build_excel_export_bytes(outputs)
        workbook = load_workbook(BytesIO(excel_bytes), data_only=True)
        self.assertIn("Consolidado", workbook.sheetnames)
        self.assertIn("Auditoria", workbook.sheetnames)
        self.assertTrue(zipfile.is_zipfile(BytesIO(excel_bytes)))
        ws = workbook["Consolidado"]
        header = [cell.value for cell in ws[1]]
        self.assertEqual("Métrica", header[0])
        self.assertEqual("Memória / fórmula", header[-1])
        pl_row = next(row for row in ws.iter_rows(min_row=2, values_only=False) if row[0].value == "PL FIDC total")
        self.assertIsInstance(pl_row[1].value, (int, float))
        self.assertNotIsInstance(pl_row[1].value, str)
        sub_row = next(row for row in ws.iter_rows(min_row=2, values_only=False) if row[0].value == "% Subordinação Total")
        self.assertAlmostEqual(0.25, sub_row[1].value, places=8)
        self.assertIn("%", sub_row[1].number_format)

        snapshot_bytes = build_consolidated_snapshot_excel_bytes(outputs)
        snapshot_workbook = load_workbook(BytesIO(snapshot_bytes), data_only=True)
        self.assertIn("Resumo exibido", snapshot_workbook.sheetnames)
        self.assertIn("Dados gráficos", snapshot_workbook.sheetnames)
        self.assertIn("Gráficos", snapshot_workbook.sheetnames)
        self.assertGreaterEqual(len(snapshot_workbook["Gráficos"]._charts), 2)

        pptx_bytes = build_pptx_export_bytes(outputs)
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        self.assertTrue(zipfile.is_zipfile(BytesIO(pptx_bytes)))
        from pptx import Presentation

        presentation = Presentation(BytesIO(pptx_bytes))
        self.assertEqual(len(outputs.fund_monthly) + 1, len(presentation.slides))
        for slide in presentation.slides:
            chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
            self.assertEqual(4, len(chart_shapes))
            self.assertTrue(all(shape.chart.has_title for shape in chart_shapes))
        with zipfile.ZipFile(BytesIO(pptx_bytes)) as archive:
            names = archive.namelist()
            xml_payload = "\n".join(
                archive.read(name).decode("utf-8", errors="ignore")
                for name in names
                if name.endswith(".xml")
            )
            chart_xml = "\n".join(
                archive.read(name).decode("utf-8")
                for name in names
                if name.startswith("ppt/charts/chart") and name.endswith(".xml")
            )
        self.assertTrue(any(name.startswith("ppt/charts/chart") for name in names))
        self.assertIn("<c:dLbls>", chart_xml)
        self.assertIn("<c:dLblPos", chart_xml)
        self.assertIn("<c:title>", chart_xml)
        self.assertIn("PL total:", xml_payload)
        self.assertIn('sz="1000"', xml_payload)

    def test_pptx_export_uses_one_2x2_slide_per_fund_plus_consolidated(self) -> None:
        dashboard_a = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard_b = _dashboard(
            fund_name="FIDC B",
            cnpj="22222222000122",
            pl_total=200.0,
            pl_senior=160.0,
            pl_mezz=20.0,
            pl_sub=20.0,
            carteira=2_000.0,
            pdd=200.0,
            buckets={4: 80.0, 7: 20.0, 8: 20.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-2",
            portfolio_name="Carteira",
            dashboards_by_cnpj={
                "11111111000111": ("FIDC A", dashboard_a),
                "22222222000122": ("FIDC B", dashboard_b),
            },
            period_label="01/2026 a 01/2026",
        )
        from pptx import Presentation

        presentation = Presentation(BytesIO(build_pptx_export_bytes(outputs)))

        self.assertEqual(3, len(presentation.slides))
        for slide in presentation.slides:
            self.assertEqual(4, sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False)))

    def test_wide_table_orders_periods_newest_to_oldest(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "11/2025",
                    "competencia_dt": pd.Timestamp("2025-11-01"),
                    "pl_total": 100.0,
                    "pl_senior": 75.0,
                    "pl_subordinada_mezz": 25.0,
                    "subordinacao_total_pct": 25.0,
                },
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "03/2026",
                    "competencia_dt": pd.Timestamp("2026-03-01"),
                    "pl_total": 110.0,
                    "pl_senior": 80.0,
                    "pl_subordinada_mezz": 30.0,
                    "subordinacao_total_pct": 27.2727,
                },
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "pl_total": 105.0,
                    "pl_senior": 78.0,
                    "pl_subordinada_mezz": 27.0,
                    "subordinacao_total_pct": 25.7143,
                },
            ]
        )
        wide = build_wide_table(monthly, scope_name="FIDC A")
        period_columns = order_period_columns_desc(wide.columns)

        self.assertEqual(["mar/26", "jan/26", "nov/25"], period_columns)
        self.assertLess(wide.columns.get_loc("mar/26"), wide.columns.get_loc("jan/26"))
        self.assertEqual("Memória / fórmula", wide.columns[-1])

    def test_consolidated_wide_html_has_section_rows_and_dense_formatting(self) -> None:
        wide = pd.DataFrame(
            [
                {
                    "Bloco": "2. PL FIDC",
                    "Métrica": "PL FIDC total",
                    "Memória / fórmula": "PATRLIQ/VL_PATRIM_LIQ",
                    "jan/26": "R$ mm 1.500,00",
                    "fev/26": "N/D",
                },
                {
                    "Bloco": "2. PL FIDC",
                    "Métrica": "% Subordinação Total",
                    "Memória / fórmula": "(Subordinada + Mezanino) / PL FIDC total.",
                    "jan/26": "27,25%",
                    "fev/26": "28,75%",
                },
                {
                    "Bloco": "7. Visão Ex-Vencidos > 360d",
                    "Métrica": "NPL Over 90d Ex 360 / Carteira Ex 360",
                    "Memória / fórmula": "NPL Over 90d Ex 360 / Carteira Ex 360.",
                    "jan/26": "12,50%",
                    "fev/26": "13,75%",
                },
                {
                    "Bloco": "7. Visão Ex-Vencidos > 360d",
                    "Métrica": "PDD / NPL Over 90d Ex 360",
                    "Memória / fórmula": "PDD Ex Over 360d / NPL Over 90d Ex 360.",
                    "jan/26": "100,00%",
                    "fev/26": "110,00%",
                },
            ]
        )

        html = _render_wide_table_html(wide)

        self.assertIn("wide-table-wrapper", html)
        self.assertIn("<details class='wide-section' style=", html)
        self.assertNotIn("<details class='wide-section' open", html)
        self.assertIn("<col class='label-col-width' style='width: 280px;'>", html)
        self.assertEqual(4, html.count("<col class='period-col-width' style='width: 96px;'>"))
        self.assertIn("<col class='formula-col-width' style='width: 340px;'>", html)
        self.assertIn("<summary>PL FIDC</summary>", html)
        self.assertIn("PL FIDC", html)
        self.assertIn("1.500,0 MM", html)
        self.assertIn("27,2%", html)
        self.assertIn("<summary>Visão Ex-Vencidos &gt; 360d</summary>", html)
        self.assertIn("<tr class='destaque'>\n<td class='label'>NPL Over 90d Ex 360 / Carteira Ex 360</td>", html)
        self.assertIn("<tr class='destaque'>\n<td class='label'>PDD / NPL Over 90d Ex 360</td>", html)
        self.assertNotIn(">N/D<", html)

    def test_dense_wide_value_uses_empty_for_missing_cells_and_br_format(self) -> None:
        self.assertEqual("", _dense_wide_value("N/D"))
        self.assertEqual("178,0%", _dense_wide_value("178,04%"))
        self.assertEqual("5.380,0 MM", _dense_wide_value("R$ mm 5.380,00"))

    def test_wide_table_uses_money_scale_by_metric_for_readability(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "scope": "consolidado",
                    "fund_name": "Carteira",
                    "cnpj": "CONSOLIDADO",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "carteira_bruta": 6_200_000_000.0,
                    "pdd_total": 320_000_000.0,
                }
            ]
        )

        wide = build_wide_table(monthly, scope_name="Carteira")
        carteira = wide.loc[wide["Métrica"] == "Carteira Bruta total", "jan/26"].iloc[0]
        pdd = wide.loc[wide["Métrica"] == "PDD total", "jan/26"].iloc[0]

        self.assertEqual("R$ mm 6.200,0", carteira)
        self.assertEqual("R$ mm 320,0", pdd)

    def test_wide_table_uses_billion_scale_only_above_one_trillion(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "scope": "consolidado",
                    "fund_name": "Carteira",
                    "cnpj": "CONSOLIDADO",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "carteira_bruta": 1_200_000_000_000.0,
                }
            ]
        )

        wide = build_wide_table(monthly, scope_name="Carteira")
        carteira = wide.loc[wide["Métrica"] == "Carteira Bruta total", "jan/26"].iloc[0]

        self.assertEqual("R$ bi 1.200,0", carteira)

    def test_wide_table_text_columns_wrap_in_fixed_grid(self) -> None:
        self.assertIn("table-layout: fixed;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("vertical-align: top;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("overflow-wrap: anywhere;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("white-space: normal;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn(".meli-kpi-grid", _MERCADO_LIVRE_UI_CSS)
        self.assertIn(".meli-kpi-card", _MERCADO_LIVRE_UI_CSS)
        td_block = _MERCADO_LIVRE_UI_CSS.split(".wide-table td {", 1)[1].split("}", 1)[0]
        self.assertNotIn("white-space: nowrap;", td_block)
        self.assertIn("white-space: normal;", td_block)
        self.assertIn("overflow-wrap: anywhere;", td_block)
        self.assertIn("line-height: 1.25;", td_block)

    def test_mercado_livre_guide_documents_usage_and_mechanics(self) -> None:
        guide = _build_mercado_livre_guide_markdown()

        self.assertIn("Como usar", guide)
        self.assertIn("Tabela Completa", guide)
        self.assertIn("Análise Crédito", guide)
        self.assertIn("Mecânica essencial", guide)
        self.assertIn("não há média simples de percentuais", guide)
        self.assertIn("NPL Over é acumulado", guide)
        self.assertIn("ex-360 remove vencidos acima de 360 dias", guide)
        self.assertIn("filtro visual", guide.lower())
        self.assertNotIn("tabelas wide", guide.lower())

    def test_display_window_full_option_keeps_loaded_36_months_by_default(self) -> None:
        available = [
            date(2023 + (idx + 3) // 12, ((idx + 3) % 12) + 1, 1)
            for idx in range(36)
        ]
        start, end = _display_window_bounds(
            selected="Todo período carregado",
            loaded_start=available[0],
            loaded_end=available[-1],
        )

        self.assertEqual(available[0], start)
        self.assertEqual(available[-1], end)

    def test_display_window_12m_filters_only_when_explicitly_selected(self) -> None:
        available = [
            date(2023 + (idx + 3) // 12, ((idx + 3) % 12) + 1, 1)
            for idx in range(36)
        ]
        start, end = _display_window_bounds(
            selected="12M",
            loaded_start=available[0],
            loaded_end=available[-1],
        )

        self.assertEqual(date(2025, 3, 1), start)
        self.assertEqual(available[-1], end)

    def test_display_window_12m_anchors_to_latest_available_when_requested_end_is_missing(self) -> None:
        available = [
            date(2024 + (idx + 1) // 12, ((idx + 1) % 12) + 1, 1)
            for idx in range(26)
            if date(2024 + (idx + 1) // 12, ((idx + 1) % 12) + 1, 1) <= date(2026, 3, 1)
        ]

        months = _display_window_months(selected="12M", available=available)

        self.assertEqual(date(2025, 3, 1), months[0])
        self.assertEqual(date(2026, 3, 1), months[-1])
        self.assertEqual(13, len(months))

    def test_requested_period_is_loaded_with_yoy_lookback_metadata(self) -> None:
        requested = build_custom_period(start_month=date(2025, 5, 1), end_month=date(2026, 4, 1))
        calculation = _period_with_yoy_lookback(requested)
        outputs = MercadoLivreOutputs(
            fund_monthly={},
            fund_wide={},
            consolidated_monthly=pd.DataFrame(),
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={},
        )

        tagged = _tag_outputs_requested_period(outputs, requested_period=requested, calculation_period=calculation)

        self.assertEqual(date(2024, 5, 1), calculation.start_month)
        self.assertEqual("12M", tagged.metadata["requested_window_option"])
        self.assertEqual("2025-05-01", tagged.metadata["requested_period_start"])
        self.assertEqual("2024-05-01", tagged.metadata["calculation_period_start"])

    def test_credit_monitor_keeps_yoy_values_after_display_filter(self) -> None:
        months = pd.date_range("2025-01-01", "2026-12-01", freq="MS")
        monthly = pd.DataFrame(
            {
                "competencia": [f"{item.month:02d}/{item.year}" for item in months],
                "competencia_dt": months,
                "fund_name": "FIDC A",
                "carteira_ex360": [100.0 + idx for idx in range(len(months))],
            }
        )
        outputs = MercadoLivreOutputs(
            fund_monthly={"1": monthly},
            fund_wide={"1": pd.DataFrame()},
            consolidated_monthly=monthly,
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={
                "display_period_months": [value.date().isoformat() for value in months[-12:]],
                "display_period_start": months[-12].date().isoformat(),
                "display_period_end": months[-1].date().isoformat(),
            },
        )

        monitor = _build_credit_monitor_for_display(outputs=outputs, display_outputs=outputs)

        self.assertEqual(12, len(monitor.consolidated_monitor))
        self.assertFalse(monitor.consolidated_monitor["carteira_ex360_yoy_pct"].isna().any())
        self.assertAlmostEqual((112.0 / 100.0 - 1.0) * 100.0, monitor.consolidated_monitor.iloc[0]["carteira_ex360_yoy_pct"])

    def test_display_window_decembers_plus_current_year_uses_non_contiguous_months(self) -> None:
        available = [
            date(2023, 11, 1),
            date(2023, 12, 1),
            date(2024, 11, 1),
            date(2025, 12, 1),
            date(2026, 1, 1),
            date(2026, 2, 1),
            date(2026, 3, 1),
            date(2026, 4, 1),
        ]

        months = _display_window_months(selected="Dezembros + Ano Atual", available=available)

        self.assertEqual(
            [
                date(2023, 12, 1),
                date(2025, 12, 1),
                date(2026, 1, 1),
                date(2026, 2, 1),
                date(2026, 3, 1),
                date(2026, 4, 1),
            ],
            months,
        )

    def test_filter_outputs_by_explicit_months_keeps_same_months_in_wide_and_metadata(self) -> None:
        monthly = pd.DataFrame(
            [
                {"competencia": "11/2023", "competencia_dt": pd.Timestamp("2023-11-01"), "fund_name": "FIDC A", "pl_total": 1.0},
                {"competencia": "12/2023", "competencia_dt": pd.Timestamp("2023-12-01"), "fund_name": "FIDC A", "pl_total": 2.0},
                {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "fund_name": "FIDC A", "pl_total": 3.0},
                {"competencia": "04/2026", "competencia_dt": pd.Timestamp("2026-04-01"), "fund_name": "FIDC A", "pl_total": 4.0},
            ]
        )
        outputs = MercadoLivreOutputs(
            fund_monthly={"1": monthly},
            fund_wide={"1": pd.DataFrame()},
            consolidated_monthly=monthly,
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={"portfolio_name": "Carteira"},
        )

        filtered = _filter_outputs_by_competencia_months(
            outputs,
            months=[date(2023, 12, 1), date(2026, 1, 1), date(2026, 4, 1)],
            label="dez/23 + jan/26 → abr/26",
            mode="Dezembros + Ano Atual",
        )

        self.assertEqual(["12/2023", "01/2026", "04/2026"], filtered.consolidated_monthly["competencia"].tolist())
        self.assertEqual(
            ["2023-12-01", "2026-01-01", "2026-04-01"],
            filtered.metadata["display_period_months"],
        )
        self.assertIn("abr/26", filtered.consolidated_wide.columns)
        self.assertIn("dez/23", filtered.consolidated_wide.columns)

    def test_snapshot_export_uses_all_displayed_months_instead_of_tail_six(self) -> None:
        months = pd.date_range("2023-12-01", "2026-04-01", freq="MS")
        monthly = pd.DataFrame(
            {
                "competencia": [f"{item.month:02d}/{item.year}" for item in months],
                "competencia_dt": months,
                "pl_total": range(len(months)),
                "pl_senior": range(len(months)),
                "pl_subordinada_mezz_ex360": range(len(months)),
                "subordinacao_total_ex360_pct": 10.0,
                "carteira_bruta": range(len(months)),
                "carteira_ex360": range(len(months)),
                "pdd_ex360": range(len(months)),
                "npl_over90_ex360": range(len(months)),
                "npl_over90_ex360_pct": 1.0,
                "pdd_npl_over90_ex360_pct": 100.0,
                "roll_rate_31_60_pct": 2.0,
            }
        )
        outputs = SimpleNamespace(consolidated_monthly=monthly)

        snapshot_bytes = build_consolidated_snapshot_excel_bytes(outputs)
        workbook = load_workbook(BytesIO(snapshot_bytes), data_only=True)
        summary_header = [cell.value for cell in workbook["Resumo exibido"][1]]
        chart_rows = list(workbook["Dados gráficos"].iter_rows(min_row=2, values_only=True))

        self.assertIn("dez/23", summary_header)
        self.assertIn("abr/26", summary_header)
        self.assertEqual(len(months), len(chart_rows))

    def test_outputs_cache_roundtrip_uses_deterministic_identity(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            loaded = load_outputs_from_cache(
                portfolio_id="outro-id-visual",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            self.assertTrue((root / "metadata.json").exists())

        self.assertIsNotNone(loaded)
        assert loaded is not None
        self.assertIn("11111111000111", loaded.fund_monthly)
        self.assertEqual(outputs.metadata["loaded_period_label"], loaded.metadata["loaded_period_label"])

    def test_outputs_cache_treats_empty_optional_warnings_as_empty_dataframe(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            (root / "warnings.csv").write_text("", encoding="utf-8")

            loaded = load_outputs_from_cache(
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )

        self.assertIsNotNone(loaded)
        assert loaded is not None
        self.assertTrue(loaded.warnings_df.empty)

    def test_outputs_cache_invalidates_empty_required_csvs(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            (root / "monthly_consolidado.csv").write_text("", encoding="utf-8")

            loaded = load_outputs_from_cache(
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )

        self.assertIsNone(loaded)

    def test_portfolio_identity_key_is_deterministic_and_duplicate_save_reuses_same_basket(self) -> None:
        funds_a = (
            PortfolioFund(cnpj="11111111000111", display_name="A"),
            PortfolioFund(cnpj="22222222000122", display_name="B"),
        )
        funds_b = (
            PortfolioFund(cnpj="22222222000122", display_name="B"),
            PortfolioFund(cnpj="11111111000111", display_name="A"),
        )

        self.assertEqual(
            portfolio_identity_key(funds_a, fallback="p1"),
            portfolio_identity_key(funds_b, fallback="p2"),
        )

        existing = PortfolioRecord(
            id="portfolio-1",
            name="Mercado Livre",
            funds=funds_a,
            created_at="2026-04-29T00:00:00Z",
            updated_at="2026-04-29T00:00:00Z",
        )
        action = _resolve_existing_portfolio_for_save(
            portfolios=[existing],
            target=None,
            name="Outro Nome",
            funds=list(funds_b),
        )

        self.assertEqual("reuse", action["action"])
        self.assertEqual(existing.id, action["portfolio"].id)


def _dashboard(
    *,
    fund_name: str,
    cnpj: str,
    pl_total: float,
    pl_senior: float,
    pl_mezz: float,
    pl_sub: float,
    carteira: float,
    pdd: float,
    buckets: dict[int, float],
) -> SimpleNamespace:
    competencia = "01/2026"
    bucket_labels = {
        1: "Até 30 dias",
        2: "31 a 60 dias",
        3: "61 a 90 dias",
        4: "91 a 120 dias",
        5: "121 a 150 dias",
        6: "151 a 180 dias",
        7: "181 a 360 dias",
        8: "361 a 720 dias",
        9: "721 a 1080 dias",
        10: "Acima de 1080 dias",
    }
    return SimpleNamespace(
        competencias=[competencia],
        fund_info={"nome_fundo": fund_name, "cnpj_fundo": cnpj, "nome_classe": "Classe única"},
        subordination_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "pl_total": pl_total,
                    "pl_senior": pl_senior,
                    "pl_mezzanino": pl_mezz,
                    "pl_subordinada_strict": pl_sub,
                    "pl_subordinada": pl_mezz + pl_sub,
                }
            ]
        ),
        default_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "direitos_creditorios": carteira,
                    "direitos_creditorios_fonte": "teste",
                    "provisao_total": pdd,
                }
            ]
        ),
        dc_canonical_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "dc_total_canonico": carteira,
                    "dc_total_fonte_efetiva": "teste",
                }
            ]
        ),
        default_buckets_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "ordem": ordem,
                    "faixa": bucket_labels[ordem],
                    "valor": valor,
                    "source_status": "reported_value",
                }
                for ordem, valor in buckets.items()
            ]
        ),
    )


if __name__ == "__main__":
    unittest.main()
